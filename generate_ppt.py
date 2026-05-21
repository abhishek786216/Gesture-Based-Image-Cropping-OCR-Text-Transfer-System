from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Plant Disease Detection and Classification"
    subtitle.text = "A Deep Learning/Machine Learning Approach\nNIT Delhi"

    # Helper function to add a standard bullet slide
    def add_bullet_slide(prs, title_text, bullet_points):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        
        if bullet_points:
            tf.text = bullet_points[0]
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
        return slide

    # Slide 2: Outline
    add_bullet_slide(prs, "Outline", [
        "Introduction",
        "Literature Review",
        "Problem Statement",
        "Motivation",
        "Objectives",
        "Proposed Framework",
        "Experimentation and Results",
        "Takeaways and Future Planning",
        "References"
    ])

    # Slide 3: Introduction
    add_bullet_slide(prs, "Introduction", [
        "What is the issue? Plant diseases significantly affect agricultural productivity and crop yield worldwide.",
        "Traditional Methods: Manual inspection by experts is time-consuming, expensive, and prone to human error.",
        "The Role of Technology: Automated systems using Computer Vision and Machine Learning offer rapid and accurate disease identification from leaf images."
    ])

    # Slide 4: Introduction (Cont.)
    add_bullet_slide(prs, "Introduction (Cont.)", [
        "Impact: Early detection of plant diseases helps in timely application of pesticides, reducing economic losses, and ensuring food security.",
        "Scope of Project: Focusing on image-based classification to distinguish between healthy and diseased plants across various species."
    ])

    # Slide 5: Literature Review
    add_bullet_slide(prs, "Literature Review", [
        "Methods heavily rely on CNNs (Convolutional Neural Networks) for image feature extraction.",
        "Wubetu Barud Demilie (2024): Provided a comparative study showing the dominance of deep learning models in precision agriculture.",
        "Gaps Identified: High computational cost, limited diverse real-world datasets, and issues with background noise."
    ])

    # Slide 6: Problem Statement
    add_bullet_slide(prs, "Problem Statement", [
        "Early and accurate detection of crop diseases is difficult due to the visual similarities between different diseases and the lack of readily available domain experts.",
        "There is a need for an automated, robust, and scalable model that can classify plant diseases from raw leaf images under varying environmental conditions."
    ])

    # Slide 7: Motivation
    add_bullet_slide(prs, "Motivation", [
        "Agricultural Importance: Agriculture is the backbone of the economy; crop protection is vital.",
        "Technological Shift: The success of AI/Deep Learning in medical imaging inspires its application in plant pathology.",
        "Empowering Farmers: Providing farmers with accessible tools (like a mobile app) to diagnose plant health instantly."
    ])

    # Slide 8: Objectives
    add_bullet_slide(prs, "Objectives", [
        "To collect and preprocess a standardized dataset of plant leaf images (e.g., PlantVillage dataset).",
        "To develop and train a Deep Learning model (like ResNet, VGG, or a custom CNN) for accurate disease classification.",
        "To compare the performance of different architectures based on accuracy, precision, and recall.",
        "To create a lightweight model suitable for edge-device deployment."
    ])

    # Slide 9: Proposed Framework
    add_bullet_slide(prs, "Proposed Framework", [
        "Data Collection: Sourcing images (Healthy vs. Diseased).",
        "Data Preprocessing: Resizing, Normalization, and Data Augmentation (rotation, flipping) to prevent overfitting.",
        "Feature Extraction & Modeling: Passing images through the chosen CNN architecture.",
        "Classification: Fully connected layers culminating in a Softmax function to predict the disease class."
    ])

    # Slide 10: Experimentation and Results
    add_bullet_slide(prs, "Experimentation and Results", [
        "Dataset Splitting: 80% Training, 10% Validation, 10% Testing.",
        "Experimental Setup: Framework used (e.g., PyTorch/TensorFlow), Hyperparameters (Learning rate, Batch size, Epochs).",
        "Evaluation Metrics: Accuracy, Precision, Recall, F1-Score."
    ])

    # Slide 11: Experimentation and Results (Cont.)
    add_bullet_slide(prs, "Experimentation and Results (Cont.)", [
        "Model Comparison:",
        "  - Model A (e.g., CNN) Accuracy: X%",
        "  - Model B (e.g., ResNet50) Accuracy: Y%",
        "Confusion Matrix: visualizes true positives vs. false positives.",
        "Loss/Accuracy Curves: demonstrates convergence during training."
    ])

    # Slide 12: Takeaways and Future Planning
    add_bullet_slide(prs, "Takeaways and Future Planning", [
        "Deep learning provides a highly reliable method for non-destructive disease detection.",
        "Pre-trained models (Transfer Learning) significantly speed up training and improve accuracy.",
        "Future Scope:",
        "  - Expanding the dataset to include multiple growth stages.",
        "  - Developing a user-friendly mobile application for farmers."
    ])

    # Slide 13: References
    add_bullet_slide(prs, "References", [
        "[1] Wubetu Barud Demilie. Plant disease detection and classification techniques: a comparative study of the performances. Journal of Big Data, 11(1):5, 2024.",
        "[2] (Add 2 or 3 other primary reference papers or datasets you used)"
    ])

    prs.save("Plant_Disease_Presentation.pptx")
    print("Presentation saved as 'Plant_Disease_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
